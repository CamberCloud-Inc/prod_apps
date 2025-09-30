import subprocess
import sys
import os
import argparse

# Install dependencies
subprocess.check_call([sys.executable, "-m", "pip", "install", "--break-system-packages", "python-docx", "python-pptx", "pypandoc"])

import docx
from pptx import Presentation
import pypandoc


def convert_docx_to_markdown(docx_path, output_path):
    """Convert Word document to Markdown using pypandoc"""
    print(f"Converting Word document using pypandoc...")

    try:
        # Use pypandoc to convert DOCX to Markdown
        output = pypandoc.convert_file(
            docx_path,
            'md',
            outputfile=output_path,
            extra_args=['--wrap=none', '--extract-media=./media']
        )
        print(f"Conversion completed successfully")
        return True
    except Exception as e:
        print(f"Error with pypandoc conversion: {str(e)}")
        print("Falling back to basic extraction...")

        # Fallback: basic extraction using python-docx
        try:
            doc = docx.Document(docx_path)
            markdown_content = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    markdown_content.append('')
                    continue

                # Check if it's a heading
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    if level.isdigit():
                        markdown_content.append(f"{'#' * int(level)} {text}")
                    else:
                        markdown_content.append(f"## {text}")
                else:
                    markdown_content.append(text)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(markdown_content))

            print("Basic extraction completed")
            return True

        except Exception as e2:
            print(f"Fallback extraction also failed: {str(e2)}")
            return False


def convert_pptx_to_markdown(pptx_path, output_path):
    """Convert PowerPoint presentation to Markdown"""
    print(f"Converting PowerPoint presentation...")

    try:
        prs = Presentation(pptx_path)
        markdown_content = []
        slide_num = 0

        for slide in prs.slides:
            slide_num += 1
            markdown_content.append(f"# Slide {slide_num}")
            markdown_content.append("")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # Check if it's a title
                    if hasattr(shape, "shape_type") and "TITLE" in str(shape.shape_type):
                        markdown_content.append(f"## {text}")
                    else:
                        # Split by lines and format as list items if multiple lines
                        lines = text.split('\n')
                        if len(lines) > 1:
                            for line in lines:
                                if line.strip():
                                    markdown_content.append(f"- {line.strip()}")
                        else:
                            markdown_content.append(text)

                    markdown_content.append("")

            markdown_content.append("---")
            markdown_content.append("")

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(markdown_content))

        print(f"Converted {slide_num} slides to Markdown")
        return True

    except Exception as e:
        print(f"Error converting PowerPoint: {str(e)}")
        return False


def main():
    parser = argparse.ArgumentParser(description='Extract Word/PowerPoint documents to Markdown')
    parser.add_argument('input_path', help='Path to the Word (.docx) or PowerPoint (.pptx) file')
    parser.add_argument('-o', '--output-dir', default='./',
                        help='Output directory for Markdown file (default: ./)')

    args = parser.parse_args()

    # Debug: print current working directory
    print(f"Current working directory: {os.getcwd()}")
    print(f"Received input_path argument: {args.input_path}")

    # Expand user path if provided
    input_path = os.path.expanduser(args.input_path)
    print(f"Expanded input_path: {input_path}")

    if not os.path.exists(input_path):
        print(f"Error: Input file not found at: {input_path}")
        sys.exit(1)

    # Create output directory if it doesn't exist
    os.makedirs(args.output_dir, exist_ok=True)

    # Determine file type
    file_ext = os.path.splitext(input_path)[1].lower()
    base_name = os.path.splitext(os.path.basename(input_path))[0]

    print(f"\nProcessing {file_ext.upper()} file: {input_path}")

    # Generate output filename
    output_filename = f"{base_name}.md"
    output_path = os.path.join(args.output_dir, output_filename)

    # Convert based on file type
    success = False

    if file_ext == '.docx':
        success = convert_docx_to_markdown(input_path, output_path)
    elif file_ext == '.pptx':
        success = convert_pptx_to_markdown(input_path, output_path)
    else:
        print(f"Error: Unsupported file type: {file_ext}")
        print("Supported formats: .docx, .pptx")
        sys.exit(1)

    if success:
        print(f"\nMarkdown file saved to: {output_path}")
        print("Office to Markdown conversion completed successfully!")
    else:
        print("\nConversion failed!")
        sys.exit(1)


if __name__ == "__main__":
    main()